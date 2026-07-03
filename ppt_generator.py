import os
import time
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import win32com.client as win32
from pathlib import Path
from datetime import datetime

# ========================= CONFIG =========================
SERPAPI_KEY = "1f693f5d4665b1344fad47eb010d0222f1ef0e94241fb63ce4a84b1b685b3431"
PEXELS_API_KEY = "B4Hp6WIVYXllRgSXBnTSbLvuUYEF8bJPBXfIvoqhV0lEagk8DT8e2ymJ"

SAVE_FOLDER = r"C:\HeartPresentations"
os.makedirs(SAVE_FOLDER, exist_ok=True)
USED_IMAGES = set()


# ========================= MENU =========================
def show_menu():
    print("\n" + "="*80)
    print(" ULTIMATE DYNAMIC PPT GENERATOR WITH ADVANCED ANIMATIONS")
    print("="*80)
    print("1. Quick Generate")
    print("2. Deep Custom")
    print("3. With File Upload")
    print("4. Full Custom")
    print("5. Exit")
    print("="*80)
    choice = input("Enter choice (1-5): ").strip()
    return choice


def get_user_input(mode):
    topic = input("\nEnter main topic: ").strip() or "Artificial Intelligence"
    custom_prompt = ""
    files = []

    if mode in ["2", "4"]:
        print("\nDescribe in detail what you want:")
        custom_prompt = input("> ").strip()

    if mode in ["3", "4"]:
        print("\nUpload files (PDF/Text):")
        while True:
            path = input("File path (or Enter to finish): ").strip()
            if not path: break
            if Path(path).exists():
                files.append(path)
            else:
                print("File not found!")

    return topic, custom_prompt, files


# ========================= READ FILES =========================
def read_user_files(file_paths):
    content = ""
    for path in file_paths:
        try:
            if path.lower().endswith('.pdf'):
                import PyPDF2
                with open(path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        content += page.extract_text() + "\n"
            else:
                with open(path, 'r', encoding='utf-8') as f:
                    content += f.read() + "\n"
        except:
            pass
    return content


# ========================= DYNAMIC RESEARCH =========================
def deep_research(topic, custom_prompt, user_content):
    print("🔍 Dynamic Research in progress...")
    research = user_content + "\n"
    if custom_prompt:
        research += f"User Focus: {custom_prompt}\n\n"

    full_text = (topic + " " + custom_prompt).lower()
    queries = [topic]

    rules = {
        "how|work|mechanism|technical": [f"{topic} working mechanism", f"{topic} technical explanation"],
        "application|use|example|case": [f"{topic} real world applications", f"{topic} case study"],
        "advantage|benefit": [f"{topic} advantages and benefits"],
        "challenge|risk|limitation": [f"{topic} challenges and limitations"],
        "future|trend|2025|2026": [f"{topic} future trends 2025-2030"],
        "market|statistic|growth": [f"{topic} market statistics"]
    }

    for key, qlist in rules.items():
        if any(k in full_text for k in key.split("|")):
            queries.extend(qlist)

    queries.extend([f"{topic} overview", f"{topic} key concepts"])

    queries = list(dict.fromkeys(queries))

    for q in queries:
        try:
            res = requests.get("https://serpapi.com/search.json",
                               params={"q": q, "api_key": SERPAPI_KEY, "num": 12}, timeout=15)
            data = res.json()
            research += " ".join([r.get("snippet", "") for r in data.get("organic_results", [])]) + ". "
        except:
            pass
        time.sleep(0.7)
    return research


# ========================= DYNAMIC SLIDES (Fixed) =========================
def generate_dynamic_slides(topic, research):
    sentences = re.split(r'[.!?]+', research)
    sentences = [re.sub(r'\s+', ' ', s).strip() for s in sentences if len(s.strip()) > 40]
    sentences = list(dict.fromkeys(sentences))

    slides = [
        {"title": topic, "bullets": [], "type": "title"},
        {"title": "Agenda", "bullets": [], "type": "agenda"}
    ]

    section_titles = ["Introduction", "Core Concepts", "Working Mechanism",
                      "Real-World Applications", "Market Statistics & Growth",
                      "Advantages & Benefits", "Challenges & Limitations",
                      "Future Scope & Trends", "Conclusion & Key Takeaways"]

    idx = 0
    for title in section_titles:
        if "Market" in title or "Statistics" in title:
            slides.append({"title": title, "bullets": [], "type": "chart"})
        elif "Conclusion" in title:
            slides.append({"title": title, "bullets": ["Major Insights", "Strategic Recommendations", "Final Outlook"], "type": "content"})
        else:
            bullets = sentences[idx:idx + 6]
            if len(bullets) < 3:
                bullets = [f"Key aspects of {title.lower()}", "Important developments", "Strategic importance"]
            slides.append({"title": title, "bullets": bullets, "type": "content"})
            idx += 6
    return slides


# ========================= DYNAMIC CHART =========================
def create_chart(topic, index):
    try:
        import matplotlib.pyplot as plt
        current_year = datetime.now().year
        years = [str(current_year - 2 + i) for i in range(6)]
        values = [45, 68, 82, 91, 97, 99]

        plt.figure(figsize=(8.8, 5.4))
        plt.bar(years, values, color=['#0088ff','#00aaff','#00ccff','#44ddff','#88eeff','#bbeeff'])
        for i, v in enumerate(values):
            plt.text(i, v + 1, f"{v}%", ha='center', color='white', fontsize=11)

        plt.title(f"{topic} Growth Projection", color='white', fontsize=15)
        plt.ylabel("Growth (%)", color='white')
        plt.gca().set_facecolor('#0f1428')
        plt.tick_params(colors='white')
        path = os.path.join(SAVE_FOLDER, f"chart_{index}.png")
        plt.savefig(path, dpi=230, bbox_inches='tight', facecolor='#0f1428')
        plt.close()
        return path
    except:
        return None


# ========================= NEXT-LEVEL DYNAMIC ANIMATIONS =========================
def apply_dynamic_animations(pptx_path):
    try:
        ppt = win32.Dispatch("PowerPoint.Application")
        ppt.Visible = False
        pres = ppt.Presentations.Open(FileName=pptx_path, ReadOnly=False, WithWindow=False)

        for slide in pres.Slides:
            slide.SlideShowTransition.EntryEffect = 54  # Morph
            slide.SlideShowTransition.Duration = 1.2

            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    if shape.TextFrame.TextRange.Text.strip():
                        anim = shape.AnimationSettings
                        anim.EntryEffect = 2      # Fade
                        anim.AdvanceMode = 2
                        anim.AdvanceTime = 0.45
                elif shape.Type == 13:  # Picture
                    try:
                        slide.TimeLine.MainSequence.AddEffect(shape, EffectId=10, Level=1, Trigger=1)  # Fly In
                        slide.TimeLine.MainSequence.AddEffect(shape, EffectId=23, Level=1, Trigger=2)  # Grow
                    except:
                        pass

        pres.Save()
        pres.Close()
        ppt.Quit()
        print(" Next-Level Dynamic Animations Applied!")
    except Exception as e:
        print(f"Animation skipped: {e}")


# ========================= CREATE PPT =========================
def create_ppt(topic, slides, theme):
    print(" Building presentation with dynamic animations...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme["bg"]

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(1.0))
        tf = title_box.text_frame
        tf.text = s["title"]
        p = tf.paragraphs[0]
        p.font.size = Pt(44 if i == 0 else 34)
        p.font.bold = True
        p.font.color.rgb = theme["text"]
        p.alignment = PP_ALIGN.CENTER

        bar = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme["accent"]

        if s.get("bullets"):
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.8), Inches(5.2))
            tf = text_box.text_frame
            tf.clear()
            tf.word_wrap = True
            for bullet in s["bullets"][:7]:
                if bullet:
                    p = tf.add_paragraph()
                    p.text = "• " + str(bullet)
                    p.font.size = Pt(19)
                    p.font.color.rgb = theme["text"]
                    p.space_after = Pt(20)

        if s.get("type") == "chart":
            chart_path = create_chart(topic, i)
            if chart_path:
                slide.shapes.add_picture(chart_path, Inches(7.7), Inches(1.8), width=Inches(5.0))
        else:
            img_url = fetch_unique_image(s["title"], topic)
            if img_url:
                img_path = download_image(img_url, f"img_{i}.jpg")
                if img_path:
                    slide.shapes.add_picture(img_path, Inches(7.85), Inches(1.8), width=Inches(4.8))

    filename = f"ULTIMATE_{topic.replace(' ', '_')}_{int(time.time())}.pptx"
    path = os.path.join(SAVE_FOLDER, filename)
    prs.save(path)

    apply_dynamic_animations(path)
    return path


def fetch_unique_image(slide_title, main_topic):
    global USED_IMAGES
    headers = {"Authorization": PEXELS_API_KEY}
    url = "https://api.pexels.com/v1/search"
    for q in [f"{slide_title} {main_topic}", f"{main_topic} {slide_title} concept"]:
        try:
            params = {"query": q, "per_page": 4, "orientation": "landscape"}
            res = requests.get(url, headers=headers, params=params, timeout=10)
            data = res.json()
            if data.get("photos"):
                for photo in data["photos"]:
                    img_url = photo["src"]["large2x"]
                    if img_url not in USED_IMAGES:
                        USED_IMAGES.add(img_url)
                        return img_url
        except:
            continue
    return None


def download_image(url, filename):
    try:
        path = os.path.join(SAVE_FOLDER, filename)
        with open(path, "wb") as f:
            f.write(requests.get(url, timeout=15).content)
        return path
    except:
        return None

# ========================= FILE EXPLORER =========================
def select_files_from_explorer():

    try:
        import tkinter as tk
        from tkinter import filedialog

        root = tk.Tk()
        root.withdraw()

        selected_files = filedialog.askopenfilenames(
            title="Select PDF or Text Files",
            filetypes=[
                ("PDF Files", "*.pdf"),
                ("Text Files", "*.txt"),
                ("All Files", "*.*")
            ]
        )

        root.destroy()

        return list(selected_files)

    except Exception as explorer_error:
        print(f" File explorer error: {explorer_error}")
        return []


# ========================= PPT GENERATION HANDLER =========================
async def handle_ppt_generation(
        topic,
        mode="1",
        custom_prompt="",
        user_files=None
):

    try:

        # ========================= DEFAULT FILES =========================

        if user_files is None:
            user_files = []

        # ========================= FILE PICKER =========================

        if mode in ["3", "4"] and not user_files:

            print("\n📂 Opening file explorer...")
            selected_files = select_files_from_explorer()

            if selected_files:
                user_files.extend(selected_files)
                print(f"Selected {len(selected_files)} file(s)")
            else:
                print("⚠️ No files selected")

        # ========================= USER CUSTOM DESCRIPTION =========================

        if mode in ["2", "4"] and not custom_prompt:

            print("\n" + "=" * 80)
            print(" DESCRIBE YOUR PRESENTATION")
            print("=" * 80)

            print("""
Describe EXACTLY what kind of presentation you want.

Examples:
- Professional business PPT
- Cybersecurity themed presentation
- Dark modern animation style
- Add statistics and graphs
- Technical deep explanation
- Short concise slides
- Investor pitch deck
- Student friendly explanation
- Add future scope and conclusion

Press ENTER twice when finished.
""")

            prompt_lines = []

            while True:

                user_line = input()

                if user_line.strip() == "":
                    break

                prompt_lines.append(user_line)

            custom_prompt = "\n".join(prompt_lines).strip()

        # ========================= THEME =========================

        ppt_theme = {
            "bg": RGBColor(15, 20, 40),
            "accent": RGBColor(0, 170, 255),
            "text": RGBColor(255, 255, 255),
            "font": "Segoe UI"
        }

        # ========================= CONSOLE INFO =========================

        print("\n" + "=" * 80)
        print(" ULTIMATE DYNAMIC PPT GENERATOR")
        print("=" * 80)
        print(f" Topic : {topic}")
        print(f"⚡ Mode  : {mode}")

        if custom_prompt:
            print(" Custom Prompt Added")

        if user_files:
            print(f"Files Attached : {len(user_files)}")

        print("=" * 80)

        # ========================= FILE CONTENT =========================

        extracted_user_content = read_user_files(user_files)

        # ========================= RESEARCH =========================

        researched_content = deep_research(
            topic,
            custom_prompt,
            extracted_user_content
        )

        # ========================= GENERATE SLIDES =========================

        generated_slides = generate_dynamic_slides(
            topic,
            researched_content
        )

        # ========================= CREATE PPT =========================

        generated_ppt_path = create_ppt(
            topic,
            generated_slides,
            ppt_theme
        )

        # ========================= OPEN PPT =========================

        try:
            os.startfile(generated_ppt_path)

        except OSError as open_error:
            print(f"⚠ Could not auto-open PPT: {open_error}")

        # ========================= SUCCESS RESPONSE =========================

        return f"""
 PRESENTATION CREATED SUCCESSFULLY

Topic:
{topic}

 Slides Generated:
{len(generated_slides)}

⚡ Mode Used:
{mode}

 Files Used:
{len(user_files)}

Saved At:
{generated_ppt_path}
"""

    except Exception as generation_error:

        return f"""
 PPT Generation Failed

Error:
{str(generation_error)}
"""
# ========================= MAIN =========================
if __name__ == "__main__":
    print("=== ULTIMATE DYNAMIC PPT GENERATOR ===\n")

    while True:
        choice = show_menu()
        if choice == "5":
            print("Goodbye!")
            break

        topic, custom_prompt, user_files = get_user_input(choice)

        user_content = read_user_files(user_files)
        research = deep_research(topic, custom_prompt, user_content)

        sentences = re.split(r'[.!?]+', research)
        sentences = [re.sub(r'\s+', ' ', s).strip() for s in sentences if len(s.strip()) > 40]
        sentences = list(dict.fromkeys(sentences))

        slides = generate_dynamic_slides(topic, research)   # Fixed: 2 arguments only

        theme = {"bg": RGBColor(15, 20, 40), "accent": RGBColor(0, 170, 255),
                 "text": RGBColor(255, 255, 255), "font": "Segoe UI"}

        ppt_path = create_ppt(topic, slides, theme)

        try:
            os.startfile(ppt_path)
        except:
            pass

        if input("\nCreate another presentation? (y/n): ").lower() != 'y':
            break